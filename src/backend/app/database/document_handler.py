"""
Document Handler
Manages PDF, PowerPoint, Word, and text documents
"""

from typing import Dict, List, Optional, Any
from pathlib import Path


class DocumentHandler:
    """Handle document file operations"""
    
    def __init__(self, base_path: Path):
        self.base_path = Path(base_path)
        self.base_path.mkdir(parents=True, exist_ok=True)
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from document"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return self._extract_pdf_text(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_powerpoint_text(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_word_text(file_path)
        elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']:
            return self._extract_text_file(file_path)
        else:
            return f"Unsupported format: {ext}"
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF"""
        try:
            # Try PyPDF2 first
            from PyPDF2 import PdfReader
            
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                text_parts.append(page.extract_text())
            
            return "\n".join(text_parts)
            
        except ImportError:
            try:
                # Fallback to pdfplumber
                import pdfplumber
                
                text_parts = []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text_parts.append(page.extract_text())
                
                return "\n".join(text_parts)
                
            except ImportError:
                return "PDF libraries not installed. Install: pip install PyPDF2 pdfplumber"
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"
    
    def _extract_powerpoint_text(self, file_path: str) -> str:
        """Extract text from PowerPoint"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
            
        except ImportError:
            return "python-pptx not installed. Install: pip install python-pptx"
        except Exception as e:
            return f"Error extracting PowerPoint: {str(e)}"
    
    def _extract_word_text(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text:
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
            
        except ImportError:
            return "python-docx not installed. Install: pip install python-docx"
        except Exception as e:
            return f"Error extracting Word: {str(e)}"
    
    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            return f"Error reading file: {str(e)}"
    
    def get_document_info(self, file_path: str) -> Dict:
        """Get document metadata"""
        ext = Path(file_path).suffix.lower()
        info = {
            "path": file_path,
            "format": ext,
            "size_mb": round(Path(file_path).stat().st_size / (1024 * 1024), 2)
        }
        
        try:
            if ext == '.pdf':
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                info["pages"] = len(reader.pages)
                info["title"] = reader.metadata.get('/Title', '') if reader.metadata else ''
                info["author"] = reader.metadata.get('/Author', '') if reader.metadata else ''
                
            elif ext in ['.pptx', '.ppt']:
                from pptx import Presentation
                prs = Presentation(file_path)
                info["slides"] = len(prs.slides)
                
            elif ext in ['.docx', '.doc']:
                from docx import Document
                doc = Document(file_path)
                info["paragraphs"] = len(doc.paragraphs)
                info["tables"] = len(doc.tables)
                
        except ImportError:
            pass
        except Exception as e:
            info["error"] = str(e)
        
        return info
    
    def search_in_document(self, file_path: str, search_term: str) -> List[Dict]:
        """Search for text in document and return occurrences"""
        text = self.extract_text(file_path)
        
        if text.startswith("Error") or text.startswith("Unsupported"):
            return [{"error": text}]
        
        # Simple line-by-line search
        results = []
        lines = text.split('\n')
        
        for line_num, line in enumerate(lines, 1):
            if search_term.lower() in line.lower():
                results.append({
                    "line": line_num,
                    "content": line.strip(),
                    "context": self._get_context(lines, line_num)
                })
        
        return results
    
    def _get_context(self, lines: List[str], line_num: int, context_lines: int = 2) -> str:
        """Get context around a line"""
        start = max(0, line_num - context_lines - 1)
        end = min(len(lines), line_num + context_lines)
        return '\n'.join(lines[start:end])
    
    def convert_to_pdf(self, file_path: str, output_path: str) -> str:
        """Convert document to PDF"""
        ext = Path(file_path).suffix.lower()
        
        try:
            if ext in ['.docx', '.doc']:
                # Use LibreOffice in production
                import subprocess
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', str(Path(output_path).parent),
                    file_path
                ], check=True)
                return output_path
                
            elif ext in ['.pptx', '.ppt']:
                # Use LibreOffice
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', str(Path(output_path).parent),
                    file_path
                ], check=True)
                return output_path
            
            else:
                return "Conversion not supported for this format"
                
        except Exception as e:
            return f"Error converting: {str(e)}"
    
    def merge_pdfs(self, file_paths: List[str], output_path: str) -> str:
        """Merge multiple PDFs into one"""
        try:
            from PyPDF2 import PdfMerger
            
            merger = PdfMerger()
            
            for pdf_path in file_paths:
                merger.append(pdf_path)
            
            merger.write(output_path)
            merger.close()
            
            return output_path
            
        except ImportError:
            return "PyPDF2 not installed. Install: pip install PyPDF2"
        except Exception as e:
            return f"Error merging PDFs: {str(e)}"
    
    def extract_images_from_pdf(self, file_path: str, output_dir: str) -> List[str]:
        """Extract images from PDF"""
        try:
            import fitz  # PyMuPDF
            from PIL import Image
            import io
            
            pdf_file = fitz.open(file_path)
            image_paths = []
            
            for page_index in range(len(pdf_file)):
                page = pdf_file[page_index]
                image_list = page.get_images()
                
                for image_index, img in enumerate(image_list, start=1):
                    xref = img[0]
                    base_image = pdf_file.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image
                    image_path = f"{output_dir}/page_{page_index + 1}_img_{image_index}.{image_ext}"
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    
                    image_paths.append(image_path)
            
            pdf_file.close()
            return image_paths
            
        except ImportError:
            return ["PyMuPDF not installed. Install: pip install PyMuPDF"]
        except Exception as e:
            return [f"Error extracting images: {str(e)}"]
    
    def create_text_summary(self, file_path: str, max_length: int = 500) -> str:
        """Create a summary of document"""
        text = self.extract_text(file_path)
        
        if text.startswith("Error"):
            return text
        
        # Simple truncation-based summary
        # In production, use NLP summarization
        if len(text) <= max_length:
            return text
        
        return text[:max_length] + "..."
